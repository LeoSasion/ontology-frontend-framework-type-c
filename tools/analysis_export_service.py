from __future__ import annotations

import argparse
import hashlib
import json
import os
import re
import shutil
import sqlite3
import tempfile
import zipfile
from datetime import datetime
from pathlib import Path
from typing import Any, Callable

from openpyxl import Workbook
from openpyxl.chart import BarChart, LineChart, PieChart, Reference
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from openpyxl.utils import get_column_letter
from openpyxl.worksheet.table import Table, TableStyleInfo

from analysis_unit_service import analysis_unit_consumer_state, get_analysis_unit, verify_analysis_unit
from query_plan_receipt_service import get_query_receipt


EXPORT_SCHEMA = "aibi-analysis-export/v1"
FIXED_ZIP_TIME = (2000, 1, 1, 0, 0, 0)
FORBIDDEN_AIBI_REPOSITORIES = {"aibi-a", "aibi-b", "aibi-d", "aibi-e", "aibi项目杂交"}
SENSITIVE_KEY_PARTS = {"password", "secret", "token", "credential", "api_key", "apikey"}
WINDOWS_PATH_PATTERN = re.compile(r"(?i)(?<![A-Za-z0-9_])([A-Za-z]:[\\/][^\s\"'<>|]+)")
SECRET_ASSIGNMENT_PATTERN = re.compile(
    r"(?i)\b(password|secret|token|credential|api[_-]?key)\b\s*[:=]\s*([^\s,;]+)"
)


def _canonical_json(value: Any) -> str:
    return json.dumps(value, ensure_ascii=False, sort_keys=True, separators=(",", ":"), default=str)


def _fingerprint(value: Any) -> str:
    return hashlib.sha256(_canonical_json(value).encode("utf-8")).hexdigest()


def _sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _redact_string(value: str) -> str:
    text = str(value)
    text = SECRET_ASSIGNMENT_PATTERN.sub(lambda match: f"{match.group(1)}=[redacted]", text)

    def redact_path(match: re.Match[str]) -> str:
        return f"[local-path-redacted]/{Path(match.group(1)).name}"

    text = WINDOWS_PATH_PATTERN.sub(redact_path, text)
    candidate = text.strip()
    if candidate and Path(candidate).is_absolute():
        return f"[local-path-redacted]/{Path(candidate).name}"
    return text


def _safe_value(value: Any, key: str = "") -> Any:
    folded = key.casefold()
    if any(part in folded for part in SENSITIVE_KEY_PARTS):
        return "[redacted]"
    if isinstance(value, dict):
        return {str(item_key): _safe_value(item, str(item_key)) for item_key, item in value.items()}
    if isinstance(value, (list, tuple)):
        return [_safe_value(item, key) for item in value]
    if isinstance(value, str):
        return _redact_string(value)
    if isinstance(value, float) and value.is_integer():
        return int(value)
    if isinstance(value, (int, float, bool)) or value is None:
        return value
    return _redact_string(str(value))


def _write_json(path: Path, value: Any) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(
        json.dumps(_safe_value(value), ensure_ascii=False, indent=2, sort_keys=True, default=str) + "\n",
        encoding="utf-8",
    )


def _reject_other_aibi_repo(path: Path) -> None:
    parts = {part.casefold() for part in path.resolve().parts}
    found = sorted(parts & FORBIDDEN_AIBI_REPOSITORIES)
    if found:
        raise ValueError(f"AIBI-C analysis export cannot write another AIBI repository: {found[0]}")


def _receipt_snapshot(receipt: dict[str, Any]) -> dict[str, Any]:
    source = receipt.get("source") if isinstance(receipt.get("source"), dict) else {}
    selection = receipt.get("selection") if isinstance(receipt.get("selection"), dict) else {}
    runtime = receipt.get("runtime") if isinstance(receipt.get("runtime"), dict) else {}
    return _safe_value({
        "schema": receipt.get("schema"),
        "receiptKey": receipt.get("receiptKey"),
        "request": receipt.get("request"),
        "status": receipt.get("status"),
        "source": {
            "tableKey": source.get("tableKey"),
            "tableKeys": source.get("tableKeys", []),
            "schemaFingerprint": source.get("schemaFingerprint"),
            "dataFingerprint": source.get("dataFingerprint"),
            "relationshipPathFingerprint": source.get("relationshipPathFingerprint"),
            "sourceFingerprint": source.get("sourceFingerprint"),
        },
        "selection": {
            "group": selection.get("group"),
            "measure": selection.get("measure"),
            "aggregation": selection.get("aggregation"),
            "filters": selection.get("filters", []),
            "joins": selection.get("joins", []),
            "semanticPlan": selection.get("semanticPlan"),
            "executionPlan": selection.get("executionPlan"),
            "knowledgeRule": selection.get("knowledgeRule"),
        },
        "runtime": {
            "engine": runtime.get("engine"),
            "executionPlanHash": runtime.get("executionPlanHash"),
            "sqlIntent": runtime.get("sqlIntent"),
        },
        "validation": receipt.get("validation", {}),
        "resultBinding": receipt.get("resultBinding"),
        "contextRefs": receipt.get("contextRefs", []),
        "domainPacks": receipt.get("domainPacks", []),
        "domainPackFingerprint": receipt.get("domainPackFingerprint"),
        "evidenceRefs": _evidence_references(receipt.get("evidenceRefs", [])),
        "unresolved": receipt.get("unresolved", []),
        "actionKey": receipt.get("actionKey"),
        "createdAt": receipt.get("createdAt"),
    })


def _evidence_references(value: Any) -> list[Any]:
    if not isinstance(value, list):
        return []
    allowed = {
        "type", "key", "receiptKey", "unitKey", "actionKey", "evidenceKey",
        "schemaFingerprint", "executionPlanHash", "resultFingerprint", "status", "engine",
    }
    references: list[Any] = []
    for item in value:
        if isinstance(item, dict):
            references.append({key: _safe_value(item.get(key), key) for key in sorted(allowed) if key in item})
        elif isinstance(item, (str, int, float, bool)) or item is None:
            references.append(_safe_value(item))
    return references


def _unit_snapshot(unit: dict[str, Any]) -> dict[str, Any]:
    return _safe_value({
        key: value
        for key, value in unit.items()
        if key not in {"workspaceId"}
    })


def _gaps(receipt: dict[str, Any], unit: dict[str, Any]) -> dict[str, Any]:
    receipt_validation = receipt.get("validation") if isinstance(receipt.get("validation"), dict) else {}
    unit_validation = unit.get("validation") if isinstance(unit.get("validation"), dict) else {}
    adapter = unit.get("chartAdapter") if isinstance(unit.get("chartAdapter"), dict) else {}
    return _safe_value({
        "unresolved": receipt.get("unresolved", []),
        "receiptBlocked": bool(receipt_validation.get("blocked")),
        "unitBlockers": unit_validation.get("blockers", []),
        "unitWarnings": unit_validation.get("warnings", []),
        "chartBlockers": adapter.get("blockers", []),
        "omissions": [
            "business database files",
            "source file contents outside the frozen Analysis Unit",
            "credentials and secrets",
            "absolute local paths",
            "compiled SQL and database paths",
            "private model reasoning",
        ],
    })


def _format_cell_value(value: Any) -> Any:
    if isinstance(value, (str, int, float, bool)) or value is None:
        return value
    return json.dumps(_safe_value(value), ensure_ascii=False, sort_keys=True)


def _autofit(worksheet, *, max_width: int = 54) -> None:
    for column_cells in worksheet.columns:
        letter = get_column_letter(column_cells[0].column)
        width = max((len(str(cell.value)) if cell.value is not None else 0 for cell in column_cells), default=0)
        worksheet.column_dimensions[letter].width = min(max(width + 2, 10), max_width)


def _style_header(cells) -> None:
    for cell in cells:
        cell.fill = PatternFill("solid", fgColor="16324F")
        cell.font = Font(color="FFFFFF", bold=True)
        cell.alignment = Alignment(vertical="center")


def _add_chart(summary, data, unit: dict[str, Any], columns: list[str]) -> None:
    adapter = unit.get("chartAdapter") if isinstance(unit.get("chartAdapter"), dict) else {}
    config = adapter.get("config") if isinstance(adapter.get("config"), dict) else {}
    chart_type = str(adapter.get("chartType") or "") if adapter.get("status") == "ready" else ""
    dimension = str(config.get("dimension") or "")
    measure = str(config.get("measure") or "")
    if chart_type not in {"bar", "line", "pie"} or dimension not in columns or measure not in columns or data.max_row < 2:
        return
    dimension_index = columns.index(dimension) + 1
    measure_index = columns.index(measure) + 1
    categories = Reference(data, min_col=dimension_index, min_row=2, max_row=data.max_row)
    values = Reference(data, min_col=measure_index, min_row=1, max_row=data.max_row)
    if chart_type == "bar":
        chart = BarChart()
        chart.type = "bar" if config.get("barOrientation") == "horizontal" else "col"
        chart.style = 10
    elif chart_type == "line":
        chart = LineChart()
        chart.style = 13
        chart.marker = "circle"
    else:
        chart = PieChart()
        chart.style = 10
    chart.add_data(values, titles_from_data=True)
    chart.set_categories(categories)
    chart.title = str(unit.get("title") or f"{measure} by {dimension}")
    chart.height = 8.0
    chart.width = 14.0
    chart.legend = None if chart_type != "pie" else chart.legend
    summary.add_chart(chart, "D3")


def _write_workbook(path: Path, receipt: dict[str, Any], unit: dict[str, Any], gaps: dict[str, Any]) -> None:
    workbook = Workbook()
    workbook.properties.creator = "AIBI-C"
    workbook.properties.lastModifiedBy = "AIBI-C"
    workbook.properties.title = str(unit.get("title") or "AIBI-C verified analysis")
    workbook.properties.subject = "Receipt-bound, frozen Analysis Unit export"
    workbook.properties.created = datetime(2000, 1, 1)
    workbook.properties.modified = datetime(2000, 1, 1)

    summary = workbook.active
    summary.title = "Summary"
    summary.sheet_view.showGridLines = False
    summary.merge_cells("A1:H1")
    summary["A1"] = "AIBI-C Verified Analysis Export"
    summary["A1"].font = Font(size=18, bold=True, color="FFFFFF")
    summary["A1"].fill = PatternFill("solid", fgColor="0F766E")
    summary["A1"].alignment = Alignment(vertical="center")
    summary.row_dimensions[1].height = 30
    adapter = unit.get("chartAdapter") if isinstance(unit.get("chartAdapter"), dict) else {}
    summary_rows = [
        ("Title", unit.get("title")),
        ("Status", unit.get("status")),
        ("Analysis kind", unit.get("kind")),
        ("Query Receipt", receipt.get("receiptKey")),
        ("Analysis Unit", unit.get("unitKey")),
        ("Result fingerprint", unit.get("resultFingerprint")),
        ("Rows", unit.get("shape", {}).get("rowCount")),
        ("Chart", adapter.get("chartType") or "blocked / table only"),
        ("Calculation", json.dumps(_safe_value(unit.get("calculation", {})), ensure_ascii=False, sort_keys=True)),
        ("Trust boundary", "Frozen scalar result only; no requery, credentials, local paths, or unrelated raw rows."),
    ]
    for row_index, (label, value) in enumerate(summary_rows, start=3):
        summary.cell(row_index, 1, label)
        summary.cell(row_index, 2, _format_cell_value(value))
        summary.cell(row_index, 1).font = Font(bold=True, color="16324F")
        summary.cell(row_index, 2).alignment = Alignment(wrap_text=True, vertical="top")
    summary.column_dimensions["A"].width = 22
    summary.column_dimensions["B"].width = 58
    summary.freeze_panes = "A3"

    data = workbook.create_sheet("Data")
    data.sheet_view.showGridLines = False
    rows = unit.get("rows") if isinstance(unit.get("rows"), list) else []
    shape = unit.get("shape") if isinstance(unit.get("shape"), dict) else {}
    columns = [str(item) for item in (shape.get("columns") or [])]
    if not columns:
        columns = sorted({str(key) for row in rows if isinstance(row, dict) for key in row})
    if columns:
        data.append(columns)
        for row in rows:
            safe_row = row if isinstance(row, dict) else {}
            data.append([_format_cell_value(_safe_value(safe_row.get(column), column)) for column in columns])
        _style_header(data[1])
        data.freeze_panes = "A2"
        data.auto_filter.ref = f"A1:{get_column_letter(len(columns))}{max(1, data.max_row)}"
        if data.max_row >= 2:
            table = Table(displayName="AnalysisData", ref=f"A1:{get_column_letter(len(columns))}{data.max_row}")
            table.tableStyleInfo = TableStyleInfo(name="TableStyleMedium2", showRowStripes=True, showColumnStripes=False)
            data.add_table(table)
        numeric = set(str(item) for item in (shape.get("numericColumns") or []))
        for column_index, column in enumerate(columns, start=1):
            if column in numeric:
                for cell in data.iter_cols(min_col=column_index, max_col=column_index, min_row=2, max_row=data.max_row):
                    for item in cell:
                        item.number_format = "#,##0.00"
        _autofit(data)
    else:
        data["A1"] = "No frozen rows"

    evidence = workbook.create_sheet("Evidence")
    evidence.sheet_view.showGridLines = False
    evidence.append(["Evidence item", "Value"])
    _style_header(evidence[1])
    evidence_rows = [
        ("Receipt status", receipt.get("status")),
        ("Source table", receipt.get("source", {}).get("tableKey")),
        ("Schema fingerprint", receipt.get("source", {}).get("schemaFingerprint")),
        ("Selection", receipt.get("selection", {})),
        ("Result binding", receipt.get("resultBinding")),
        ("Evidence references", receipt.get("evidenceRefs", [])),
        ("Unit grain", unit.get("grain", {})),
        ("Unit validation", unit.get("validation", {})),
        ("Chart adapter", adapter),
        ("Gaps and omissions", gaps),
    ]
    for label, value in evidence_rows:
        evidence.append([label, json.dumps(_safe_value(value), ensure_ascii=False, sort_keys=True)])
    evidence.freeze_panes = "A2"
    evidence.column_dimensions["A"].width = 24
    evidence.column_dimensions["B"].width = 70
    for row in evidence.iter_rows(min_row=2, min_col=2, max_col=2):
        row[0].alignment = Alignment(wrap_text=True, vertical="top")

    manifest = workbook.create_sheet("Manifest")
    manifest.sheet_view.showGridLines = False
    manifest.append(["Contract", "Value"])
    _style_header(manifest[1])
    manifest_rows = [
        ("Schema", EXPORT_SCHEMA),
        ("Projection", "scalar-result-fields-only"),
        ("Receipt fingerprint", receipt.get("resultBinding", {}).get("resultFingerprint")),
        ("Unit fingerprint", unit.get("resultFingerprint")),
        ("Definition fingerprint", unit.get("definitionFingerprint")),
        ("Requery", "No"),
        ("Business database write", "No"),
        ("Credential inclusion", "No"),
        ("Absolute path inclusion", "No"),
    ]
    for label, value in manifest_rows:
        manifest.append([label, _format_cell_value(value)])
    manifest.column_dimensions["A"].width = 28
    manifest.column_dimensions["B"].width = 72

    thin = Side(style="thin", color="D6E2EA")
    for sheet in workbook.worksheets:
        used = sheet.calculate_dimension()
        for row in sheet[used]:
            for cell in row:
                if cell.value is not None:
                    cell.border = Border(bottom=thin)
                    cell.alignment = Alignment(
                        horizontal=cell.alignment.horizontal,
                        vertical=cell.alignment.vertical or "top",
                        wrap_text=cell.alignment.wrap_text,
                    )
    _add_chart(summary, data, unit, columns)
    workbook.save(path)
    _normalize_zip(path)


def _write_docx(path: Path, receipt: dict[str, Any], unit: dict[str, Any], gaps: dict[str, Any]) -> None:
    try:
        from docx import Document
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.shared import Pt, RGBColor
    except ImportError as error:
        raise RuntimeError("DOCX export requires python-docx; install requirements.txt before requesting docx") from error
    document = Document()
    document.core_properties.author = "AIBI-C"
    document.core_properties.last_modified_by = "AIBI-C"
    document.core_properties.title = str(unit.get("title") or "AIBI-C verified analysis")
    document.core_properties.subject = "Receipt-bound, frozen Analysis Unit export"
    normal = document.styles["Normal"]
    normal.font.name = "Microsoft YaHei"
    normal.font.size = Pt(10.5)
    title = document.add_heading(str(unit.get("title") or "AIBI-C 可验证分析报告"), level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.LEFT
    for run in title.runs:
        run.font.color.rgb = RGBColor(15, 118, 110)
    document.add_paragraph("AIBI-C · Receipt-bound verified analysis", style="Subtitle")
    summary = document.add_table(rows=0, cols=2)
    summary.style = "Light Shading Accent 1"
    for label, value in (
        ("Status", unit.get("status")),
        ("Analysis kind", unit.get("kind")),
        ("Query Receipt", receipt.get("receiptKey")),
        ("Analysis Unit", unit.get("unitKey")),
        ("Result fingerprint", unit.get("resultFingerprint")),
        ("Rows", (unit.get("shape") or {}).get("rowCount")),
    ):
        cells = summary.add_row().cells
        cells[0].text = str(label)
        formatted = _format_cell_value(value)
        cells[1].text = str("—" if formatted is None else formatted)
    document.add_heading("口径与计算 / Grain and calculation", level=1)
    document.add_paragraph(json.dumps(_safe_value(unit.get("grain", {})), ensure_ascii=False, indent=2, sort_keys=True))
    document.add_paragraph(json.dumps(_safe_value(unit.get("calculation", {})), ensure_ascii=False, indent=2, sort_keys=True))
    rows = unit.get("rows") if isinstance(unit.get("rows"), list) else []
    shape = unit.get("shape") if isinstance(unit.get("shape"), dict) else {}
    columns = [str(item) for item in (shape.get("columns") or [])]
    if not columns:
        columns = sorted({str(key) for row in rows if isinstance(row, dict) for key in row})
    document.add_heading("冻结结果 / Frozen result", level=1)
    if columns and rows:
        table = document.add_table(rows=1, cols=len(columns))
        table.style = "Light Shading Accent 1"
        for index, column in enumerate(columns):
            table.rows[0].cells[index].text = column
        for row in rows[:100]:
            cells = table.add_row().cells
            record = row if isinstance(row, dict) else {}
            for index, column in enumerate(columns):
                formatted = _format_cell_value(_safe_value(record.get(column), column))
                cells[index].text = str("" if formatted is None else formatted)
        if len(rows) > 100:
            document.add_paragraph(f"Only the first 100 of {len(rows)} frozen rows are displayed in DOCX; the complete frozen payload remains in the export bundle.")
    else:
        document.add_paragraph("No frozen rows.")
    document.add_heading("验证与限制 / Validation and limitations", level=1)
    validation = unit.get("validation") if isinstance(unit.get("validation"), dict) else {}
    for label, items in (("Blockers", validation.get("blockers", [])), ("Warnings", validation.get("warnings", [])), ("Omissions", gaps.get("omissions", []))):
        document.add_paragraph(f"{label}: {', '.join(map(str, items)) or 'None'}")
    document.add_paragraph("Trust boundary: generated only from the stored Query Receipt and frozen Analysis Unit. No requery and no business database write.")
    document.save(path)
    _normalize_zip(path)


def _write_pptx(path: Path, receipt: dict[str, Any], unit: dict[str, Any], gaps: dict[str, Any]) -> None:
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt
    except ImportError as error:
        raise RuntimeError("PPTX export requires python-pptx; install requirements.txt before requesting pptx") from error
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    presentation.core_properties.author = "AIBI-C"
    presentation.core_properties.last_modified_by = "AIBI-C"
    presentation.core_properties.title = str(unit.get("title") or "AIBI-C verified analysis")

    def add_title(slide, title: str, subtitle: str = "") -> None:
        slide.shapes.title.text = title
        paragraph = slide.shapes.title.text_frame.paragraphs[0]
        paragraph.font.name = "Microsoft YaHei"
        paragraph.font.size = Pt(28)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(15, 118, 110)
        if subtitle and len(slide.placeholders) > 1:
            slide.placeholders[1].text = subtitle

    cover = presentation.slides.add_slide(presentation.slide_layouts[0])
    add_title(cover, str(unit.get("title") or "AIBI-C 可验证分析"), f"Receipt {receipt.get('receiptKey')} · Unit {unit.get('unitKey')}")

    overview = presentation.slides.add_slide(presentation.slide_layouts[1])
    add_title(overview, "结论与可信状态 / Conclusion and trust")
    overview.placeholders[1].text = "\n".join([
        f"Status: {unit.get('status') or '—'}",
        f"Analysis kind: {unit.get('kind') or '—'}",
        f"Frozen rows: {(unit.get('shape') or {}).get('rowCount') or 0}",
        f"Result fingerprint: {str(unit.get('resultFingerprint') or '')[:20]}",
        "Source: stored Query Receipt + frozen Analysis Unit",
    ])

    data_slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    add_title(data_slide, "冻结结果摘要 / Frozen result snapshot")
    rows = unit.get("rows") if isinstance(unit.get("rows"), list) else []
    shape = unit.get("shape") if isinstance(unit.get("shape"), dict) else {}
    columns = [str(item) for item in (shape.get("columns") or [])]
    if not columns:
        columns = sorted({str(key) for row in rows if isinstance(row, dict) for key in row})
    visible_columns = columns[:6]
    visible_rows = rows[:8]
    if visible_columns:
        table = data_slide.shapes.add_table(len(visible_rows) + 1, len(visible_columns), Inches(0.6), Inches(1.5), Inches(12.1), Inches(4.8)).table
        for column_index, column in enumerate(visible_columns):
            table.cell(0, column_index).text = column
        for row_index, row in enumerate(visible_rows, start=1):
            record = row if isinstance(row, dict) else {}
            for column_index, column in enumerate(visible_columns):
                formatted = _format_cell_value(_safe_value(record.get(column), column))
                table.cell(row_index, column_index).text = str("" if formatted is None else formatted)
        for cell in table.iter_cells():
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.name = "Microsoft YaHei"
                paragraph.font.size = Pt(10)
                paragraph.alignment = PP_ALIGN.LEFT
    else:
        textbox = data_slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(11), Inches(1))
        textbox.text_frame.text = "No frozen rows."

    evidence = presentation.slides.add_slide(presentation.slide_layouts[1])
    add_title(evidence, "证据与限制 / Evidence and limitations")
    validation = unit.get("validation") if isinstance(unit.get("validation"), dict) else {}
    evidence.placeholders[1].text = "\n".join([
        f"Receipt status: {receipt.get('status') or '—'}",
        f"Blockers: {', '.join(map(str, validation.get('blockers', []))) or 'None'}",
        f"Warnings: {', '.join(map(str, validation.get('warnings', []))) or 'None'}",
        f"Omissions: {', '.join(map(str, gaps.get('omissions', []))) or 'None'}",
        "No requery · No business database write · No credentials or local paths",
    ])
    presentation.save(path)
    _normalize_zip(path)


def _normalize_zip(path: Path) -> None:
    normalized = path.with_suffix(path.suffix + ".normalized")
    with zipfile.ZipFile(path, "r") as source, zipfile.ZipFile(normalized, "w", compression=zipfile.ZIP_DEFLATED) as target:
        for name in sorted(source.namelist()):
            payload = source.read(name)
            if name == "docProps/core.xml":
                for tag in (b"created", b"modified"):
                    pattern = rb"(<dcterms:" + tag + rb"[^>]*>)[^<]*(</dcterms:" + tag + rb">)"
                    payload = re.sub(pattern, rb"\g<1>2000-01-01T00:00:00Z\g<2>", payload)
            info = zipfile.ZipInfo(name, FIXED_ZIP_TIME)
            info.compress_type = zipfile.ZIP_DEFLATED
            info.external_attr = 0o100644 << 16
            target.writestr(info, payload)
    os.replace(normalized, path)


def _write_deterministic_archive(staging: Path, archive_path: Path) -> None:
    temporary = archive_path.with_suffix(archive_path.suffix + ".tmp")
    with zipfile.ZipFile(temporary, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        for path in sorted(item for item in staging.rglob("*") if item.is_file()):
            relative = path.relative_to(staging).as_posix()
            info = zipfile.ZipInfo(relative, FIXED_ZIP_TIME)
            info.compress_type = zipfile.ZIP_DEFLATED
            info.external_attr = 0o100644 << 16
            archive.writestr(info, path.read_bytes())
    os.replace(temporary, archive_path)


def _report_markdown(receipt: dict[str, Any], unit: dict[str, Any], gaps: dict[str, Any]) -> str:
    adapter = unit.get("chartAdapter") if isinstance(unit.get("chartAdapter"), dict) else {}
    validation = unit.get("validation") if isinstance(unit.get("validation"), dict) else {}
    lines = [
        "# AIBI-C 可验证分析报告",
        "",
        f"- 结论：{_redact_string(str(unit.get('title') or '-'))}",
        f"- 状态：{unit.get('status') or '-'}",
        f"- 分析类型：{unit.get('kind') or '-'}",
        f"- Query Receipt：{receipt.get('receiptKey') or '-'}",
        f"- Analysis Unit：{unit.get('unitKey') or '-'}",
        f"- 结果指纹：{unit.get('resultFingerprint') or '-'}",
        f"- 图表：{adapter.get('chartType') or '未适配'}",
        "",
        "## 口径与粒度",
        "",
        f"```json\n{json.dumps(_safe_value(unit.get('grain', {})), ensure_ascii=False, indent=2, sort_keys=True)}\n```",
        "",
        "## 已验证计算",
        "",
        f"```json\n{json.dumps(_safe_value(unit.get('calculation', {})), ensure_ascii=False, indent=2, sort_keys=True)}\n```",
        "",
        "## 验证状态",
        "",
        f"- Blockers：{', '.join(map(str, validation.get('blockers', []))) or '无'}",
        f"- Warnings：{', '.join(map(str, validation.get('warnings', []))) or '无'}",
        f"- Chart blockers：{', '.join(map(str, adapter.get('blockers', []))) or '无'}",
        "",
        "## 缺口与省略",
        "",
    ]
    omissions = gaps.get("omissions") if isinstance(gaps.get("omissions"), list) else []
    lines.extend(f"- {item}" for item in omissions)
    lines.extend([
        "",
        "## 可信边界",
        "",
        "本报告只消费已存 Query Receipt、冻结 Analysis Unit、Chart Adapter 与证据引用；导出阶段不重新查询，不写业务数据库。",
        "",
    ])
    return "\n".join(lines)


def export_analysis_command(
    args: argparse.Namespace,
    *,
    open_db: Callable[[], Any],
    active_workspace_id: Callable[[sqlite3.Connection], str],
    root: Path,
) -> dict[str, Any]:
    receipt_key = str(args.receipt or "").strip()
    unit_key = str(args.unit or "").strip()
    if not receipt_key or not unit_key:
        raise ValueError("receipt and unit are required")
    requested_formats = list(dict.fromkeys(str(item).strip().casefold() for item in (getattr(args, "format", None) or []) if str(item).strip()))
    formats = requested_formats or ["xlsx", "md"]
    unsupported_formats = sorted(set(formats) - {"xlsx", "docx", "pptx", "md"})
    if unsupported_formats:
        raise ValueError(f"Unsupported analysis export format: {unsupported_formats[0]}")
    with open_db() as connection:
        workspace_id = active_workspace_id(connection)
        receipt = get_query_receipt(connection, workspace_id, receipt_key)
        unit = get_analysis_unit(connection, workspace_id, unit_key)
        consumer_state = analysis_unit_consumer_state(connection, workspace_id, unit) if unit else None
    if not receipt:
        raise ValueError(f"Unknown query receipt in active workspace: {receipt_key}")
    if not unit:
        raise ValueError(f"Unknown Analysis Unit in active workspace: {unit_key}")
    if str(unit.get("queryReceiptKey") or "") != receipt_key:
        raise ValueError("Analysis Unit is not bound to the requested Query Receipt")
    if not consumer_state or not consumer_state.get("usable"):
        blockers = ", ".join((consumer_state or {}).get("blockers") or ["analysis-unit-source-drifted"])
        raise ValueError(f"Analysis export rejects a drifted Analysis Unit: {blockers}")
    result_binding = receipt.get("resultBinding") if isinstance(receipt.get("resultBinding"), dict) else {}
    if str(receipt.get("status") or "") != "executed" or str(unit.get("status") or "") != "ready":
        raise ValueError("Analysis export requires an executed Query Receipt and ready Analysis Unit")
    if str(result_binding.get("resultFingerprint") or "") != str(unit.get("resultFingerprint") or ""):
        raise ValueError("Analysis Unit result fingerprint does not match the Query Receipt")
    verification = verify_analysis_unit(unit)
    verification["freshness"] = consumer_state
    if not verification.get("ok"):
        raise ValueError("Analysis Unit recomputation verification failed")

    export_root_value = str(os.environ.get("AIBI_EXPORT_ROOT") or "").strip()
    export_root = Path(export_root_value).expanduser() if export_root_value else root / "data" / "local" / "exports"
    _reject_other_aibi_repo(export_root)
    export_root.mkdir(parents=True, exist_ok=True)
    output_value = str(args.output or "").strip()
    archive_path = Path(output_value).expanduser() if output_value else export_root / f"aibi-analysis-{unit_key}.zip"
    if not archive_path.is_absolute():
        archive_path = export_root / archive_path
    if archive_path.suffix.casefold() != ".zip":
        archive_path = archive_path.with_suffix(".zip")
    _reject_other_aibi_repo(archive_path)
    archive_path.parent.mkdir(parents=True, exist_ok=True)

    staging = Path(tempfile.mkdtemp(prefix="aibi-analysis-export-", dir=export_root))
    try:
        safe_receipt = _receipt_snapshot(receipt)
        safe_unit = _unit_snapshot(unit)
        gaps = _gaps(safe_receipt, safe_unit)
        _write_json(staging / "snapshots" / "query-receipt.json", safe_receipt)
        _write_json(staging / "snapshots" / "analysis-unit.json", safe_unit)
        _write_json(staging / "snapshots" / "chart-adapter.json", safe_unit.get("chartAdapter", {}))
        _write_json(staging / "gaps.json", gaps)
        if "md" in formats:
            report = _report_markdown(safe_receipt, safe_unit, gaps)
            (staging / "report.md").write_text(report, encoding="utf-8")
        if "xlsx" in formats:
            _write_workbook(staging / "result.xlsx", safe_receipt, safe_unit, gaps)
        if "docx" in formats:
            _write_docx(staging / "result.docx", safe_receipt, safe_unit, gaps)
        if "pptx" in formats:
            _write_pptx(staging / "result.pptx", safe_receipt, safe_unit, gaps)

        files = sorted(path for path in staging.rglob("*") if path.is_file())
        file_entries = [
            {"path": path.relative_to(staging).as_posix(), "sizeBytes": path.stat().st_size, "sha256": _sha256(path)}
            for path in files
        ]
        manifest = {
            "schema": EXPORT_SCHEMA,
            "receiptKey": receipt_key,
            "unitKey": unit_key,
            "workspace": hashlib.sha256(workspace_id.encode("utf-8")).hexdigest()[:16],
            "status": "ready",
            "projection": "scalar-result-fields-only",
            "formats": formats,
            "sourceCreatedAt": safe_receipt.get("createdAt"),
            "files": file_entries,
            "contentFingerprint": _fingerprint(file_entries),
            "gaps": gaps,
            "exclusions": gaps["omissions"],
            "sideEffects": {"requery": False, "businessDatabaseWrite": False},
        }
        _write_json(staging / "manifest.json", manifest)
        _write_deterministic_archive(staging, archive_path)
    finally:
        shutil.rmtree(staging, ignore_errors=True)

    return {
        "ok": True,
        "analysisExport": {
            "schema": EXPORT_SCHEMA,
            "archivePath": str(archive_path),
            "archiveSha256": _sha256(archive_path),
            "receiptKey": receipt_key,
            "unitKey": unit_key,
            "fileCount": len(manifest["files"]) + 1,
            "formats": formats,
            "contentFingerprint": manifest["contentFingerprint"],
            "manifest": manifest,
            "verification": verification,
        },
        "artifacts": [{"label": "analysis-export", "path": str(archive_path), "sha256": _sha256(archive_path)}],
        "evidence": [
            {"type": "queryPlanReceipt", "receiptKey": receipt_key},
            {"type": "analysisUnit", "unitKey": unit_key, "resultFingerprint": unit.get("resultFingerprint")},
        ],
    }
