from __future__ import annotations

import hashlib
import json
import os
import shutil
import sqlite3
import subprocess
import tempfile
import zipfile
from pathlib import Path
from typing import Any

from openpyxl import load_workbook
from docx import Document
from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
VERIFY_DIR = Path(tempfile.mkdtemp(prefix="aibi-analysis-export-"))
KEEP_VERIFY_DIR = os.environ.get("AIBI_KEEP_ANALYSIS_EXPORT", "").strip() == "1"
DB_PATH = VERIFY_DIR / "runtime.sqlite"
ENV = {
    **os.environ,
    "AIBI_HYBRID_DB_PATH": str(DB_PATH),
    "AIBI_HYBRID_DUCKDB_PATH": str(VERIFY_DIR / "runtime.duckdb"),
    "AIBI_EVIDENCE_BUNDLE_ROOT": str(VERIFY_DIR / "evidence"),
    "AIBI_EXPORT_ROOT": str(VERIFY_DIR / "exports"),
    "PYTHONIOENCODING": "utf-8",
}
checks: list[dict[str, Any]] = []


def check(label: str, ok: bool, detail: Any = None) -> None:
    checks.append({"label": label, "ok": bool(ok), "detail": detail})


def run(label: str, arguments: list[str], expected_status: int = 0) -> dict[str, Any]:
    completed = subprocess.run(
        ["python", "tools/aibi_cli.py", "--json", *arguments],
        cwd=ROOT,
        env=ENV,
        text=True,
        encoding="utf-8",
        capture_output=True,
        check=False,
    )
    try:
        parsed = json.loads(completed.stdout.strip())
    except json.JSONDecodeError:
        parsed = None
    ok = completed.returncode == expected_status and (expected_status != 0 or parsed and parsed.get("ok") is True)
    check(label, ok, {
        "status": completed.returncode,
        "stderr": completed.stderr[-1000:],
        "stdout": completed.stdout[-1000:] if not ok else "",
    })
    return {"status": completed.returncode, "parsed": parsed, "stderr": completed.stderr, "stdout": completed.stdout}


def database_counts() -> dict[str, int]:
    with sqlite3.connect(DB_PATH) as connection:
        return {
            table: int(connection.execute(f"SELECT COUNT(*) FROM {table}").fetchone()[0])
            for table in ["table_registry", "query_plan_receipts", "analysis_units", "action_drafts"]
        }


def sha256_bytes(value: bytes) -> str:
    return hashlib.sha256(value).hexdigest()


try:
    run("import", ["import-commit", "validation-inputs/orders.csv", "--table", "orders", "--name", "Orders", "--mode", "create", "--yes"])
    comparison = run("comparison-query", [
        "query", "--table", "orders", "--group", "channel", "--measure", "net_sales", "--agg", "sum",
        "--request", "各渠道净销售额柱状图",
    ])["parsed"] or {}
    receipt = comparison.get("queryPlanReceipt") or {}
    unit = comparison.get("analysisUnit") or {}
    check(
        "query-produced-ready-receipt-bound-unit",
        receipt.get("status") == "executed"
        and unit.get("status") == "ready"
        and unit.get("queryReceiptKey") == receipt.get("receiptKey"),
    )

    secret_value = "M10_SUPERSECRET"
    rebuilt = run("redaction-fixture-unit", [
        "analysis-unit-build", "--receipt", str(receipt.get("receiptKey")), "--kind", "comparison",
        "--rows-json", json.dumps(comparison.get("rows") or [], ensure_ascii=False),
        "--title", f"渠道分析 token={secret_value} C:\\Users\\Example\\private.csv",
        "--preferred-chart", "bar",
    ])["parsed"] or {}
    unit = rebuilt.get("analysisUnit") or unit
    before_counts = database_counts()

    archive_a = VERIFY_DIR / "exports" / "analysis-a.zip"
    archive_b = VERIFY_DIR / "exports" / "analysis-b.zip"
    export_a = run("first-export", [
        "export-analysis", "--receipt", str(receipt.get("receiptKey")), "--unit", str(unit.get("unitKey")),
        "--output", str(archive_a),
    ])["parsed"] or {}
    export_b = run("repeat-export", [
        "export-analysis", "--receipt", str(receipt.get("receiptKey")), "--unit", str(unit.get("unitKey")),
        "--output", str(archive_b),
    ])["parsed"] or {}
    after_counts = database_counts()

    bytes_a = archive_a.read_bytes() if archive_a.exists() else b""
    bytes_b = archive_b.read_bytes() if archive_b.exists() else b""
    archive_hashes = {"a": sha256_bytes(bytes_a), "b": sha256_bytes(bytes_b)}
    differing_entries: dict[str, list[str]] = {}
    if bytes_a and bytes_b and bytes_a != bytes_b:
        with zipfile.ZipFile(archive_a, "r") as first, zipfile.ZipFile(archive_b, "r") as second_archive:
            differing_entries["bundle"] = [
                name for name in sorted(set(first.namelist()) | set(second_archive.namelist()))
                if name not in first.namelist() or name not in second_archive.namelist() or first.read(name) != second_archive.read(name)
            ]
            if "result.xlsx" in first.namelist() and "result.xlsx" in second_archive.namelist():
                first_xlsx = VERIFY_DIR / "first-debug.xlsx"
                second_xlsx = VERIFY_DIR / "second-debug.xlsx"
                first_xlsx.write_bytes(first.read("result.xlsx"))
                second_xlsx.write_bytes(second_archive.read("result.xlsx"))
                with zipfile.ZipFile(first_xlsx, "r") as first_book, zipfile.ZipFile(second_xlsx, "r") as second_book:
                    differing_entries["xlsx"] = [
                        name for name in sorted(set(first_book.namelist()) | set(second_book.namelist()))
                        if name not in first_book.namelist() or name not in second_book.namelist() or first_book.read(name) != second_book.read(name)
                    ]
    check(
        "same-receipt-and-unit-produce-identical-archive",
        bool(bytes_a) and bytes_a == bytes_b
        and export_a.get("analysisExport", {}).get("archiveSha256") == sha256_bytes(bytes_a)
        and export_b.get("analysisExport", {}).get("archiveSha256") == sha256_bytes(bytes_b),
        {"hashes": archive_hashes, "differingEntries": differing_entries},
    )
    check("export-does-not-write-business-database", before_counts == after_counts, {"before": before_counts, "after": after_counts})

    expected_files = {
        "gaps.json",
        "manifest.json",
        "report.md",
        "result.xlsx",
        "snapshots/analysis-unit.json",
        "snapshots/chart-adapter.json",
        "snapshots/query-receipt.json",
    }
    with zipfile.ZipFile(archive_a, "r") as archive:
        names = set(archive.namelist())
        manifest = json.loads(archive.read("manifest.json").decode("utf-8"))
        entries = {item["path"]: item for item in manifest.get("files", [])}
        hashes_match = all(
            name in entries and entries[name]["sha256"] == sha256_bytes(archive.read(name))
            for name in expected_files - {"manifest.json"}
        )
        check("bundle-has-exact-auditable-files", names == expected_files, sorted(names))
        check(
            "manifest-hashes-and-side-effects-are-explicit",
            hashes_match
            and manifest.get("schema") == "aibi-analysis-export/v1"
            and manifest.get("projection") == "scalar-result-fields-only"
            and manifest.get("sideEffects") == {"businessDatabaseWrite": False, "requery": False},
        )
        text_payload = b"\n".join(
            archive.read(name)
            for name in names
            if name.endswith((".json", ".md"))
        ).decode("utf-8", errors="ignore")
        xlsx_bytes = archive.read("result.xlsx")

    xlsx_path = VERIFY_DIR / "result.xlsx"
    xlsx_path.write_bytes(xlsx_bytes)
    with zipfile.ZipFile(xlsx_path, "r") as workbook_archive:
        workbook_payload = b"\n".join(workbook_archive.read(name) for name in workbook_archive.namelist())
    combined = text_payload.encode("utf-8") + workbook_payload
    check(
        "export-redacts-secrets-and-absolute-paths",
        secret_value.encode("utf-8") not in combined
        and str(VERIFY_DIR).encode("utf-8") not in combined
        and b"C:\\Users\\Example\\private.csv" not in combined
        and b"compiledSql" not in combined
        and b"[redacted]" in combined
        and b"[local-path-redacted]/private.csv" in combined,
    )

    workbook = load_workbook(xlsx_path, data_only=False)
    check("workbook-has-structured-sheets", workbook.sheetnames == ["Summary", "Data", "Evidence", "Manifest"], workbook.sheetnames)
    data = workbook["Data"]
    exported_rows = [
        {str(data.cell(1, column).value): data.cell(row, column).value for column in range(1, data.max_column + 1)}
        for row in range(2, data.max_row + 1)
    ]
    check(
        "workbook-preserves-only-frozen-scalar-result",
        exported_rows == unit.get("rows")
        and all(not isinstance(value, (dict, list)) for row in exported_rows for value in row.values()),
    )
    check("workbook-contains-native-adapter-chart", len(workbook["Summary"]._charts) == 1)
    formulas = [
        cell.value
        for sheet in workbook.worksheets
        for row in sheet.iter_rows()
        for cell in row
        if isinstance(cell.value, str) and cell.value.startswith("=")
    ]
    check("workbook-does-not-introduce-recalculation-formulas", formulas == [], formulas)

    office_path = VERIFY_DIR / "exports" / "analysis-office.zip"
    office = run("office-export", [
        "export-analysis", "--receipt", str(receipt.get("receiptKey")), "--unit", str(unit.get("unitKey")),
        "--format", "docx", "--format", "pptx", "--output", str(office_path),
    ])["parsed"] or {}
    with zipfile.ZipFile(office_path, "r") as office_archive:
        office_names = set(office_archive.namelist())
        docx_bytes = office_archive.read("result.docx")
        pptx_bytes = office_archive.read("result.pptx")
        office_manifest = json.loads(office_archive.read("manifest.json").decode("utf-8"))
    docx_path = VERIFY_DIR / "result.docx"
    pptx_path = VERIFY_DIR / "result.pptx"
    docx_path.write_bytes(docx_bytes)
    pptx_path.write_bytes(pptx_bytes)
    document = Document(docx_path)
    presentation = Presentation(pptx_path)
    with zipfile.ZipFile(docx_path, "r") as document_archive, zipfile.ZipFile(pptx_path, "r") as presentation_archive:
        office_xml = b"\n".join(document_archive.read(name) for name in document_archive.namelist()) + b"\n" + b"\n".join(presentation_archive.read(name) for name in presentation_archive.namelist())
    check(
        "docx-and-pptx-are-native-verified-artifacts",
        office_names == {"gaps.json", "manifest.json", "result.docx", "result.pptx", "snapshots/analysis-unit.json", "snapshots/chart-adapter.json", "snapshots/query-receipt.json"}
        and len(document.paragraphs) > 4
        and len(presentation.slides) == 4
        and office_manifest.get("formats") == ["docx", "pptx"]
        and office.get("analysisExport", {}).get("formats") == ["docx", "pptx"],
        {"names": sorted(office_names), "paragraphs": len(document.paragraphs), "slides": len(presentation.slides)},
    )
    check(
        "office-artifacts-preserve-the-same-redaction-boundary",
        secret_value.encode("utf-8") not in office_xml
        and b"C:\\Users\\Example\\private.csv" not in office_xml
        and str(VERIFY_DIR).encode("utf-8") not in office_xml,
    )

    second = run("second-query", [
        "query", "--table", "orders", "--measure", "net_sales", "--agg", "sum", "--request", "净销售额指标卡",
    ])["parsed"] or {}
    mismatch_path = VERIFY_DIR / "exports" / "mismatch.zip"
    mismatch = run("mismatched-receipt-unit-blocked", [
        "export-analysis", "--receipt", str(receipt.get("receiptKey")),
        "--unit", str((second.get("analysisUnit") or {}).get("unitKey")), "--output", str(mismatch_path),
    ], expected_status=1)
    check(
        "mismatched-receipt-unit-produces-no-artifact",
        not mismatch_path.exists() and "not bound" in str((mismatch.get("parsed") or {}).get("error", "")),
    )

    forbidden_path = VERIFY_DIR / "AIBI-D" / "forbidden.zip"
    forbidden = run("other-aibi-output-blocked", [
        "export-analysis", "--receipt", str(receipt.get("receiptKey")), "--unit", str(unit.get("unitKey")),
        "--output", str(forbidden_path),
    ], expected_status=1)
    check(
        "other-aibi-output-is-rejected-before-write",
        not forbidden_path.exists() and "another AIBI repository" in str((forbidden.get("parsed") or {}).get("error", "")),
    )

    cli_contract = run("cli-contract", ["cli-contract"])["parsed"] or {}
    commands = (cli_contract.get("contract") or {}).get("commands") or []
    export_contract = next((item for item in commands if item.get("name") == "export-analysis"), {})
    check(
        "export-capability-is-artifact-only",
        export_contract.get("mutationMode") == "artifact-export"
        and export_contract.get("writesBusinessState") is False
        and export_contract.get("writesEvidence") is False
        and export_contract.get("writesArtifacts") is True,
        export_contract,
    )

finally:
    failed = [item for item in checks if not item["ok"]]
    print(json.dumps({
        "ok": not failed,
        "schema": "aibi-analysis-export-verify/v1",
        "generatedBy": "scripts/verify-analysis-export.py",
        "verifyDir": str(VERIFY_DIR) if KEEP_VERIFY_DIR else None,
        "checks": checks,
        "failedChecks": failed,
    }, ensure_ascii=False, indent=2))
    if not KEEP_VERIFY_DIR:
        shutil.rmtree(VERIFY_DIR, ignore_errors=True)

raise SystemExit(1 if failed else 0)
